import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
import os

def pdf_to_pptx(pdf_path, pptx_path, dpi=300):
    # 打开 PDF 文档
    doc = fitz.open(pdf_path)
    # 创建一个新的 PPTX 演示文稿
    presentation = Presentation()

    # 设置幻灯片尺寸为 16:9
    presentation.slide_width = Inches(13.33)
    presentation.slide_height = Inches(7.5)

    for i in range(len(doc)):
        # 选择一个空白幻灯片布局
        slide_layout = presentation.slide_layouts[5]
        slide = presentation.slides.add_slide(slide_layout)
        
        # 将 PDF 页面转换为高分辨率图片
        page = doc.load_page(i)
        zoom = dpi / 72  # PDF 默认分辨率是 72 DPI
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat)
        img_path = f"page_{i}.png"
        pix.save(img_path)
        
        # 计算保持比例的宽高
        img_width, img_height = pix.width, pix.height
        aspect_ratio = img_width / img_height

        # 计算图片放入幻灯片后的尺寸和位置
        slide_width = presentation.slide_width
        slide_height = presentation.slide_height
        if aspect_ratio > (slide_width / slide_height):
            # 图片宽度占满幻灯片宽度
            width = slide_width
            height = width / aspect_ratio
            top = (slide_height - height) / 2
            left = 0
        else:
            # 图片高度占满幻灯片高度
            height = slide_height
            width = height * aspect_ratio
            left = (slide_width - width) / 2
            top = 0
        
        # 将图片添加到幻灯片中
        slide.shapes.add_picture(img_path, left, top, width=width, height=height)
        
        # 删除临时图片文件
        os.remove(img_path)

    # 保存 PPTX 文件
    presentation.save(pptx_path)

# 路径设置
pdf_path = "5-19.pdf"
pptx_path = "presentation.pptx"

# 转换 PDF 为 PPTX
pdf_to_pptx(pdf_path, pptx_path, dpi=300)  # 使用 300 DPI 的分辨率
